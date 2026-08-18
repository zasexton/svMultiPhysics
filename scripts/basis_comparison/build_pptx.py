"""Build a PowerPoint deck with the basis-comparison figures embedded as
true SVG images (PNG fallback for compatibility).

Pipeline:
  1. Scaffold a .pptx using python-pptx with PNGs for layout.
  2. Post-process the .pptx ZIP to inject the corresponding SVG files and
     attach the OOXML 2016+ svgBlip extension to each picture so PowerPoint
     renders the SVG natively (no rasterization).
"""

from __future__ import annotations

import io
import re
import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# --------------------------------------------------------------------------
# Slide order — matches the presentation narrative.
# --------------------------------------------------------------------------
SLIDE_ORDER = [
    # === ACCURACY ===
    # Group A1 — internal correctness
    "kronecker_delta",
    "partition_of_unity_diagnostic",
    "polynomial_reproduction",
    # Group A2 — direct comparison
    "reference_element_contours",
    "cross_sections",
    "scatter_legacy_vs_oop",
    # Group A3 — permutation
    "node_permutation_diagram",
    # Group A4 — floating-point agreement
    "ulp_distance_histogram",
    "error_cdf",
    # Group A5 — derivatives
    "gradient_field_quiver",
    # Group A6 — derived FE quantities
    "mass_matrix_comparison",
    "conditioning_bars",
    # Group A7 — end-to-end
    "interpolation_error_agreement",
    "interpolation_error_cdfs",
    "summary_max_error_bars",
    # === PERFORMANCE ===
    "perf_microbench_pointwise",
    "perf_setup_cost",
    "perf_steady_state_access",
    "perf_fused_kernels",
    "perf_cache_effectiveness",
    "perf_parallel_scaling",
]

# OOXML namespace and IDs for the SVG extension (PowerPoint 2016+).
SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
SVG_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"


def scaffold_pptx(figures_dir: Path, output: Path) -> list[str]:
    """Create the initial .pptx with PNG images.

    Returns the slide names in order so we can match them up during
    post-processing.
    """
    prs = Presentation()
    # Widescreen 16:9.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank
    slide_names: list[str] = []

    for name in SLIDE_ORDER:
        png = figures_dir / f"{name}.png"
        if not png.exists():
            print(f"  SKIP {name}: missing {png.name}")
            continue
        slide = prs.slides.add_slide(blank_layout)
        # Centered image filling the slide with a small margin.
        margin = Inches(0.25)
        slide_w = prs.slide_width - 2 * margin
        slide_h = prs.slide_height - 2 * margin
        # Pick the dimension that fits the image's aspect ratio without crop.
        # python-pptx will preserve aspect when only width or height is set.
        slide.shapes.add_picture(str(png), margin, margin,
                                 width=slide_w, height=slide_h)
        slide_names.append(name)

    prs.save(str(output))
    return slide_names


def inject_svg(pptx_path: Path, figures_dir: Path, slide_names: list[str]) -> None:
    """Post-process the .pptx to embed real SVG via svgBlip extension.

    For each slide:
      - Add ppt/media/svg{N}.svg next to the existing PNG.
      - Add a new relationship in the slide's _rels referencing the SVG.
      - Edit the slide XML to add an a:extLst/a:ext/asvg:svgBlip child to the
        existing a:blip element so PowerPoint uses the SVG as the source.
    """
    tmp_path = pptx_path.with_suffix(".pptx.tmp")
    with zipfile.ZipFile(pptx_path, "r") as src, \
         zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as dst:

        # Collect existing Content_Types and patch in image/svg+xml if missing.
        ct_data = src.read("[Content_Types].xml").decode("utf-8")
        if 'Extension="svg"' not in ct_data:
            ct_data = ct_data.replace(
                "</Types>",
                '<Default Extension="svg" ContentType="image/svg+xml"/></Types>',
            )

        # Track which slide-xml files we've patched for SVG.
        for item in src.namelist():
            data = src.read(item)

            if item == "[Content_Types].xml":
                dst.writestr(item, ct_data)
                continue

            # Match each slide XML.
            slide_match = re.match(r"ppt/slides/slide(\d+)\.xml$", item)
            if slide_match:
                slide_idx = int(slide_match.group(1)) - 1  # 0-based
                if slide_idx < len(slide_names):
                    name = slide_names[slide_idx]
                    svg_path = figures_dir / f"{name}.svg"
                    if svg_path.exists():
                        data = patch_slide_xml(data, slide_idx)
                        # Add SVG file under ppt/media/.
                        media_name = f"ppt/media/svg{slide_idx + 1}.svg"
                        dst.writestr(media_name, svg_path.read_bytes())
                        # Patch the slide's _rels file.
                        rels_name = f"ppt/slides/_rels/slide{slide_idx + 1}.xml.rels"
                        rels_data = src.read(rels_name).decode("utf-8")
                        rels_data = patch_slide_rels(rels_data, slide_idx)
                        dst.writestr(rels_name, rels_data)
                dst.writestr(item, data)
                continue

            # Skip slide rels here — they're written in the slide branch above.
            if re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels$", item):
                m = re.match(r"ppt/slides/_rels/slide(\d+)\.xml\.rels$", item)
                idx = int(m.group(1)) - 1
                if idx < len(slide_names):
                    name = slide_names[idx]
                    if (figures_dir / f"{name}.svg").exists():
                        # Already written as part of the slide branch.
                        continue
                # Otherwise pass through unchanged.
                dst.writestr(item, data)
                continue

            dst.writestr(item, data)

    shutil.move(str(tmp_path), str(pptx_path))


def patch_slide_xml(xml_bytes: bytes, slide_idx: int) -> bytes:
    """Inject the asvg:svgBlip extension into the existing a:blip element.

    Existing structure (from python-pptx):
        <a:blip r:embed="rIdN"/>
    Becomes:
        <a:blip r:embed="rIdN">
          <a:extLst>
            <a:ext uri="{96DAC541-7B7A-43D3-8B79-37D633B846F1}">
              <asvg:svgBlip
                  xmlns:asvg="http://schemas.microsoft.com/office/drawing/2016/SVG/main"
                  r:embed="rIdSVG"/>
            </a:ext>
          </a:extLst>
        </a:blip>

    The SVG relationship ID we'll use is one greater than the highest existing
    rId in the slide rels (rels patcher allocates rIdSVG accordingly).
    """
    xml_str = xml_bytes.decode("utf-8")

    # Find the existing a:blip self-closing tag.
    # python-pptx emits something like: <a:blip r:embed="rId4"/>
    # We replace it with an opened blip + extension + close.
    blip_re = re.compile(r'<a:blip\s+r:embed="(rId\d+)"\s*/>')
    m = blip_re.search(xml_str)
    if m is None:
        return xml_bytes  # no picture on this slide

    existing_rid = m.group(1)
    # Compute SVG relId by finding all rIds in the file and adding one.
    all_rids = [int(x) for x in re.findall(r'r:embed="rId(\d+)"', xml_str)]
    next_rid = max(all_rids) + 1 if all_rids else 1
    svg_rid = f"rId{next_rid}"

    replacement = (
        f'<a:blip r:embed="{existing_rid}">'
        f'<a:extLst>'
        f'<a:ext uri="{SVG_EXT_URI}">'
        f'<asvg:svgBlip xmlns:asvg="{ASVG_NS}" r:embed="{svg_rid}"/>'
        f'</a:ext>'
        f'</a:extLst>'
        f'</a:blip>'
    )
    xml_str = blip_re.sub(replacement, xml_str, count=1)
    return xml_str.encode("utf-8")


def patch_slide_rels(rels_str: str, slide_idx: int) -> str:
    """Add a relationship for the SVG image alongside the existing PNG."""
    # Find max existing rId.
    rids = [int(x) for x in re.findall(r'Id="rId(\d+)"', rels_str)]
    next_rid = max(rids) + 1 if rids else 1
    svg_rid = f"rId{next_rid}"
    target = f"../media/svg{slide_idx + 1}.svg"
    new_rel = (
        f'<Relationship Id="{svg_rid}" '
        f'Type="{SVG_REL_TYPE}" '
        f'Target="{target}"/>'
    )
    return rels_str.replace("</Relationships>", new_rel + "</Relationships>")


def main() -> int:
    here = Path(__file__).resolve().parent
    figures_dir = here / "figures"
    output = here / "basis_comparison.pptx"

    print("Scaffolding deck with PNG fallbacks...")
    slide_names = scaffold_pptx(figures_dir, output)
    print(f"  scaffolded {len(slide_names)} slides: {output}")

    print("Injecting true SVG via svgBlip extension...")
    inject_svg(output, figures_dir, slide_names)

    size_mb = output.stat().st_size / (1024 * 1024)
    print(f"Wrote {output} ({size_mb:.1f} MB, {len(slide_names)} slides).")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
