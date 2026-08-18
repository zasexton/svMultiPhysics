#!/usr/bin/env python3
"""
Generate a short slide deck covering the FE Basis library capabilities.

Produces: BasisModule_Overview.pptx  (in the same directory as this script)

Requires: python-pptx, matplotlib, numpy
"""

import os
import sys
import importlib.util

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_PPTX   = os.path.join(SCRIPT_DIR, "BasisModule_Overview.pptx")
PNG_DIR    = os.path.join(SCRIPT_DIR, ".basis_png_tmp")

# ── Colours ─────────────────────────────────────────────────────────────

DARK   = RGBColor(0x2C, 0x3E, 0x50)
BLUE   = RGBColor(0x3A, 0x7C, 0xA5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xEC, 0xF0, 0xF1)
ACCENT = RGBColor(0xD4, 0x62, 0x2B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════════════════════
# Step 1 — Re-render figures as high-res PNGs
# ═══════════════════════════════════════════════════════════════════════════

def render_pngs():
    """Import the figure module and re-save each figure as PNG."""
    os.makedirs(PNG_DIR, exist_ok=True)

    spec = importlib.util.spec_from_file_location(
        "gen_figs", os.path.join(SCRIPT_DIR, "generate_basis_figures.py"))
    mod = importlib.util.module_from_spec(spec)

    # Patch OUT_DIR so the module writes PNGs to our tmp dir
    original_out = None

    def _patch():
        nonlocal original_out
        original_out = mod.OUT_DIR
        mod.OUT_DIR = PNG_DIR

    spec.loader.exec_module(mod)
    _patch()

    # Also patch each figure function to save PNG instead of SVG
    funcs = [
        ("fig_feature_coverage",        "feature_coverage"),
        ("fig_dof_scaling",             "dof_scaling"),
        ("fig_hessian_support",         "hessian_support"),
        ("fig_polynomial_reproduction", "polynomial_reproduction"),
        ("fig_basis_family_breadth",    "basis_family_breadth"),
    ]

    # Monkey-patch plt.Figure.savefig to force PNG + higher DPI
    _orig_savefig = plt.Figure.savefig
    def _png_savefig(self, fname, *args, **kwargs):
        if isinstance(fname, str) and fname.endswith(".svg"):
            fname = fname[:-4] + ".png"
        kwargs["format"] = "png"
        kwargs["dpi"] = 200
        return _orig_savefig(self, fname, *args, **kwargs)

    plt.Figure.savefig = _png_savefig
    try:
        for func_name, _ in funcs:
            getattr(mod, func_name)()
    finally:
        plt.Figure.savefig = _orig_savefig

    paths = {}
    for _, name in funcs:
        p = os.path.join(PNG_DIR, f"{name}.png")
        if os.path.exists(p):
            paths[name] = p
        else:
            print(f"  WARNING: {p} not generated")
    return paths


# ═══════════════════════════════════════════════════════════════════════════
# Step 2 — Slide-building helpers
# ═══════════════════════════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, *,
                    font_size=16, color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def add_figure(slide, png_path, left, top, width, height=None):
    """Add a PNG image scaled to fit within width (and optionally height)."""
    if height:
        slide.shapes.add_picture(png_path, left, top, width, height)
    else:
        slide.shapes.add_picture(png_path, left, top, width=width)


def add_accent_bar(slide, left, top, width, height, color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_subtitle_bar(slide, top, text, *, bg_color=BLUE, font_color=WHITE):
    """Full-width coloured subtitle strip."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, Inches(0.55))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = font_color
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER


# ═══════════════════════════════════════════════════════════════════════════
# Step 3 — Build each slide
# ═══════════════════════════════════════════════════════════════════════════

def build_deck(png_paths):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank = prs.slide_layouts[6]  # blank layout

    # ── Slide 1: Title ──────────────────────────────────────────────────

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, DARK)
    add_accent_bar(slide, Inches(0), Inches(2.7), SLIDE_W, Inches(0.08))

    add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                "FE Basis Module", font_size=44, bold=True, color=WHITE)
    add_textbox(slide, Inches(1), Inches(2.85), Inches(11), Inches(0.8),
                "Capabilities Overview & Comparison with Legacy Solver",
                font_size=24, color=LIGHT)
    add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
                "svMultiPhysics  |  Code/Source/solver/FE/Basis/",
                font_size=16, color=RGBColor(0x95, 0xA5, 0xA6))

    # ── Slide 2: Architecture Overview ──────────────────────────────────

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                "Architecture Overview", font_size=32, bold=True, color=DARK)

    left_items = [
        "Reference-space only: all coordinates in parametric space; "
        "physical mapping handled externally by the Geometry module",
        "Analytic derivatives preferred: Lagrange, Spectral, Hermite, "
        "Hierarchical, Bernstein, B-Spline, and tensor-product bases "
        "provide analytic gradients and Hessians (no finite differences)",
        "Semantic caching via BasisCache keyed on evaluation-relevant "
        "state (knots, weights, order, element type)",
        "SoA memory layout for SIMD-friendly access in BasisCacheEntry "
        "and BatchedBasisData",
        "VTK-compatible node ordering for seamless visualization output",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5),
                    left_items, font_size=14, spacing=Pt(10))

    # Right side: architecture text block
    arch_text = (
        "BasisRequest  -->  basis_factory::create()\n"
        "                        |\n"
        "  +----- Scalar --------+------- Vector -----+\n"
        "  |                                           |\n"
        "  |-- LagrangeBasis          RaviartThomasBasis\n"
        "  |-- HierarchicalBasis      NedelecBasis\n"
        "  |-- BernsteinBasis         BDMBasis\n"
        "  |-- SpectralBasis (GLL)\n"
        "  |-- SerendipityBasis\n"
        "  |-- HermiteBasis\n"
        "  |-- BubbleBasis\n"
        "  |-- BSplineBasis\n"
        "  |-- TensorProductBasis\n"
        "  +-- NURBSTensorBasis"
    )
    txBox = slide.shapes.add_textbox(
        Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = arch_text
    p.font.size = Pt(12)
    p.font.color.rgb = DARK
    p.font.name = "Consolas"
    p.alignment = PP_ALIGN.LEFT

    # Surrounding box for the arch text
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.6), Inches(1.15), Inches(6.3), Inches(5.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFA)
    box.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
    box.line.width = Pt(1)
    # Send box behind text
    sp = box._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

    # ── Slide 3: Feature Coverage ───────────────────────────────────────

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                "Element Topology x Basis Family Support",
                font_size=32, bold=True, color=DARK)

    if "feature_coverage" in png_paths:
        add_figure(slide, png_paths["feature_coverage"],
                   Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5))

    add_textbox(slide, Inches(0.8), Inches(6.85), Inches(11), Inches(0.5),
                "Legacy solver supports only Lagrange basis on 14 fixed element types.  "
                "New library adds 12 additional families including H(div) and H(curl) vector spaces.",
                font_size=13, color=RGBColor(0x7F, 0x8C, 0x8D))

    # ── Slide 4: DOF Scaling ────────────────────────────────────────────

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                "DOF Scaling: Arbitrary Order vs Fixed Order",
                font_size=32, bold=True, color=DARK)

    if "dof_scaling" in png_paths:
        add_figure(slide, png_paths["dof_scaling"],
                   Inches(1.5), Inches(1.2), Inches(10.0), Inches(5.6))

    add_textbox(slide, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.5),
                "Legacy solver is locked to p = 1 and p = 2 (black circles).  "
                "New library supports arbitrary polynomial order on all 7 canonical topologies.",
                font_size=13, color=RGBColor(0x7F, 0x8C, 0x8D))

    # ── Slide 5: Derivative / Hessian Support ───────────────────────────

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                "Derivative Support by Element Type",
                font_size=32, bold=True, color=DARK)

    if "hessian_support" in png_paths:
        add_figure(slide, png_paths["hessian_support"],
                   Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.0))

    bullet_items = [
        "Legacy solver provides analytic Hessians for only 4 of 14 element "
        "types (Tri6, Quad8, Quad9, Tet10)",
        "New Basis library provides analytic values, gradients, and Hessians "
        "for all element types across all basis families",
        "Hessians are required for stabilization terms (e.g. SUPG, VMS) -- "
        "legacy falls back to finite differences for unsupported elements",
    ]
    add_bullet_list(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(1.3),
                    bullet_items, font_size=12,
                    color=RGBColor(0x7F, 0x8C, 0x8D), spacing=Pt(4))

    # ── Slide 6: Polynomial Reproduction ────────────────────────────────

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                "Polynomial Reproduction Accuracy",
                font_size=32, bold=True, color=DARK)

    if "polynomial_reproduction" in png_paths:
        add_figure(slide, png_paths["polynomial_reproduction"],
                   Inches(0.5), Inches(1.1), Inches(12.0), Inches(5.5))

    add_textbox(slide, Inches(0.8), Inches(6.75), Inches(11.5), Inches(0.6),
                "Nodal interpolation through Lagrange basis reproduces all monomials "
                "up to basis order to machine precision (~1e-14 to 1e-16).  "
                "Verified across Line, Triangle, Quad, Tetrahedra, and Hexahedra for p = 1 and p = 2.",
                font_size=13, color=RGBColor(0x7F, 0x8C, 0x8D))

    # ── Slide 7: Capability Breadth ─────────────────────────────────────

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)
    add_accent_bar(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                "Basis Library Capability Comparison",
                font_size=32, bold=True, color=DARK)

    if "basis_family_breadth" in png_paths:
        add_figure(slide, png_paths["basis_family_breadth"],
                   Inches(1.5), Inches(1.2), Inches(10.0), Inches(5.2))

    add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.6),
                "13x more basis families  |  4x more element-order combinations  |  "
                "3 functional spaces (H1, H(div), H(curl)) vs 1  |  "
                "386 unit tests vs 0",
                font_size=14, bold=True, color=BLUE)

    # ── Slide 8: Summary ────────────────────────────────────────────────

    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, DARK)
    add_accent_bar(slide, Inches(0), Inches(1.3), SLIDE_W, Inches(0.06))

    add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.8),
                "Key Takeaways", font_size=36, bold=True, color=WHITE)

    summary = [
        "13 basis families vs 1 -- Lagrange, Hierarchical, Bernstein, "
        "Spectral/GLL, Serendipity, Hermite, Bubble, B-Spline, NURBS, "
        "Raviart-Thomas, Nedelec, BDM",

        "Arbitrary polynomial order on all 7 canonical topologies "
        "(Line, Triangle, Quad, Tetra, Hex, Wedge, Pyramid) -- "
        "legacy is fixed at p = 1 and p = 2",

        "Full analytic Hessians for every element type -- "
        "critical for VMS/SUPG stabilization without finite-difference fallback",

        "H(div) and H(curl) vector bases enable mixed formulations "
        "(Stokes, Darcy, Maxwell) not possible with the legacy solver",

        "386 unit tests across 16 test files verify partition of unity, "
        "polynomial completeness, derivative accuracy, and error paths",

        "Reference-space-only design with BasisCache and BatchEvaluator "
        "for efficient integration into the assembly pipeline",
    ]

    tf = add_bullet_list(slide, Inches(1), Inches(1.7), Inches(11), Inches(5.0),
                         summary, font_size=17, color=LIGHT, spacing=Pt(14))

    prs.save(OUT_PPTX)
    print(f"\nDeck saved to {OUT_PPTX}")


# ═══════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("Rendering figures as PNG...")
    png_paths = render_pngs()
    print(f"\n{len(png_paths)} figures ready.  Building deck...")
    build_deck(png_paths)
    print("Done.")
